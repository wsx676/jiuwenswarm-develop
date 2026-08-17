#!/usr/bin/env python3
"""Generate references/template.pptx — the neutral blank deck template.

    python3 scripts/make_blank_template.py [-o references/template.pptx]

The template is deliberately plain: PowerPoint's own default theme, no logo, no
background art, no confidentiality chrome. It exists to give merge_slides.py a
real master/layout structure to inherit, and to supply the cover and closing
pages that finalize_deck.py's default --order "t1,s*,t5" pulls in.

Structure (matching what finalize_deck.py documents as template slide roles):

    t1  cover    Title Slide layout; ctrTitle left EMPTY for fill_cover.py
    t2  TOC      Title Only layout
    t3  blank    Blank layout
    t4  blank    Blank layout
    t5  closing  Blank layout, centered "Thank You"

Two structural details are load-bearing and must not be "simplified" away:

  - The cover's ctrTitle paragraph carries an explicit <a:pPr>. python-pptx
    emits a bare <a:p/> for an untouched placeholder, and fill_cover.py needs a
    paragraph with properties to replace runs inside.
  - The cover's subtitle placeholder ships three paragraphs already labelled
    部门：/ 汇报人：/ 日期：, because fill_cover.py --meta appends values after
    existing label text rather than creating lines.

Regenerating this file is expected to be reproducible: the output is a function
of python-pptx's bundled default template plus the edits below.
"""
import logging
import sys
from argparse import ArgumentParser
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Program output goes to stdout, diagnostics to stderr. Both travel through
# logging, with a bare "%(message)s" format so the text is unchanged and the two
# streams stay separate for anything parsing this tool's output.
_OUT = logging.getLogger("make_blank_template.out")
_OUT.propagate = False
_OUT.setLevel(logging.INFO)
_out_handler = logging.StreamHandler(sys.stdout)
_out_handler.setFormatter(logging.Formatter("%(message)s"))
_OUT.addHandler(_out_handler)

LOGGER = logging.getLogger("make_blank_template")
LOGGER.propagate = False
LOGGER.setLevel(logging.INFO)
_err_handler = logging.StreamHandler(sys.stderr)
_err_handler.setFormatter(logging.Formatter("%(message)s"))
LOGGER.addHandler(_err_handler)


def emit(line):
    """Program output on stdout."""
    _OUT.info(line)

SKILL_ROOT = Path(__file__).resolve().parent.parent

# Office theme Accent1. Kept in sync with THEME.accent in scripts/components.js.
ACCENT = RGBColor(0x44, 0x72, 0xC4)
TEXT = RGBColor(0x1D, 0x1D, 0x1A)
CJK_FONT = "Microsoft YaHei"

# python-pptx's default template layout order. Index -> role we use it for.
LAYOUT_TITLE = 0   # "Title Slide"      -> slideLayout1.xml
LAYOUT_TITLE_ONLY = 5  # "Title Only"   -> slideLayout6.xml
LAYOUT_BLANK = 6   # "Blank"            -> slideLayout7.xml

COVER_META_LABELS = ("部门：", "汇报人：", "日期：")


def style(run, size, color=TEXT, bold=False):
    run.font.name = CJK_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def build_cover(pres):
    """Cover with an EMPTY title placeholder — fill_cover.py writes it later.

    The title box is moved off the default centered position onto the same
    left-aligned axis addOpeningSlide draws, so a deck built through the
    template path and one built purely in PptxGenJS look like the same family.
    """
    slide = pres.slides.add_slide(pres.slide_layouts[LAYOUT_TITLE])
    title = slide.placeholders[0]
    title.left, title.top = Inches(1.11), Inches(1.00)
    title.width, title.height = Inches(8.56), Inches(0.75)

    # Give the empty paragraph explicit properties. Without a <a:pPr> the
    # placeholder serializes as <a:p/> and fill_cover.py cannot find a
    # paragraph to write into.
    #
    # The alignment is load-bearing, not cosmetic: the Title Slide layout
    # centers its placeholders, so an unset paragraph puts the filled-in title
    # in the middle of the slide while the accent rule below stays hard left.
    # Setting .alignment creates the <a:pPr> on its own -- verified to emit the
    # same XML as reaching into the private element to pre-create it.
    para = title.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    para.font.name = CJK_FONT
    para.font.size = Pt(29)
    para.font.bold = True
    para.font.color.rgb = ACCENT

    # Accent rule under the title, mirroring addOpeningSlide's marker.
    bar = slide.shapes.add_shape(1, Inches(1.11), Inches(1.90), Inches(0.79), Inches(0.023))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    subtitle = slide.placeholders[1]
    subtitle.left, subtitle.top = Inches(1.11), Inches(2.20)
    subtitle.width, subtitle.height = Inches(4.0), Inches(1.20)
    frame = subtitle.text_frame
    frame.word_wrap = True
    for i, label in enumerate(COVER_META_LABELS):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        style(para.add_run(), 14, TEXT)
        para.runs[0].text = label
    return slide


def build_toc(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[LAYOUT_TITLE_ONLY])
    title = slide.placeholders[0]
    title.left, title.top = Inches(1.13), Inches(0.40)
    title.width, title.height = Inches(5.33), Inches(0.67)
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "目录"
    style(run, 29, TEXT, bold=True)
    return slide


def build_blank(pres):
    return pres.slides.add_slide(pres.slide_layouts[LAYOUT_BLANK])


def build_closing(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[LAYOUT_BLANK])
    box = slide.shapes.add_textbox(Inches(0), Inches(3.20), Inches(13.333), Inches(1.00))
    para = box.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = "Thank You"
    style(run, 36, ACCENT, bold=True)
    return slide


def main():
    ap = ArgumentParser(description=__doc__)
    ap.add_argument("-o", "--out", default=str(SKILL_ROOT / "references" / "template.pptx"))
    args = ap.parse_args()

    pres = Presentation()
    pres.slide_width, pres.slide_height = Inches(13.333), Inches(7.5)

    build_cover(pres)
    build_toc(pres)
    build_blank(pres)
    build_blank(pres)
    build_closing(pres)

    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    pres.save(out)
    emit(f"make_blank_template: wrote {out} — {len(pres.slides)} slides, "
          f"{len(pres.slide_masters)} master(s), {len(pres.slide_layouts)} layouts, 16:9")


if __name__ == "__main__":
    main()
